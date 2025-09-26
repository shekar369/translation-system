"""
Parser Worker - Extracts text content from various document formats
"""
import asyncio
import logging
import uuid
import tempfile
import os
from sqlalchemy.orm import Session
from app.models.database import get_db
from app.models.jobs import JobFile, JobArtifact, FileStatus
from app.services.queue_service import QueueService, JobEvents
from app.services.storage_service import StorageService

logger = logging.getLogger(__name__)

class ParserWorker:
    def __init__(self):
        self.queue_service = QueueService()
        self.storage_service = StorageService()
        self.db = next(get_db())

    async def start(self):
        """Start the parser worker"""
        logger.info("Starting Parser Worker")
        await self.queue_service.consume("parsing", self.handle_parse_request)

    async def handle_parse_request(self, message: dict):
        """Handle file parsing request"""
        job_id = message.get("job_id")
        file_id = message.get("file_id")
        object_key = message.get("object_key")
        media_type = message.get("media_type")
        mime_type = message.get("mime_type")

        logger.info(f"Parsing file {file_id} for job {job_id}")

        try:
            # Update file status
            job_file = self.db.query(JobFile).filter(JobFile.id == uuid.UUID(file_id)).first()
            if job_file:
                job_file.status = FileStatus.PARSING
                self.db.commit()

            # Parse the file based on media type
            if media_type == "document":
                parsed_content = await self.parse_document(object_key, mime_type)
            elif media_type == "image":
                parsed_content = await self.parse_image(object_key)
            else:
                # For audio/video, we'll extract metadata only
                parsed_content = await self.extract_media_metadata(object_key, mime_type)

            # Save parsed content to storage
            parsed_object_key = f"parsed/{job_id}/{file_id}.json"
            await self.save_parsed_content(parsed_content, parsed_object_key)

            # Create artifact record
            artifact = JobArtifact(
                job_id=uuid.UUID(job_id),
                job_file_id=uuid.UUID(file_id),
                artifact_type="parsed",
                object_key=parsed_object_key
            )
            self.db.add(artifact)

            # Update file status
            if job_file:
                job_file.status = FileStatus.COMPLETED
                self.db.commit()

            # Notify orchestrator
            response = {
                "event": JobEvents.PARSING_COMPLETED,
                "job_id": job_id,
                "file_id": file_id,
                "success": True,
                "parsed_object_key": parsed_object_key
            }
            await self.queue_service.publish("events", response)

            logger.info(f"Successfully parsed file {file_id}")

        except Exception as e:
            logger.error(f"Error parsing file {file_id}: {e}")

            # Update file status to failed
            if job_file:
                job_file.status = FileStatus.FAILED
                self.db.commit()

            # Notify orchestrator of failure
            response = {
                "event": JobEvents.PARSING_FAILED,
                "job_id": job_id,
                "file_id": file_id,
                "success": False,
                "error": str(e)
            }
            await self.queue_service.publish("events", response)

    async def parse_document(self, object_key: str, mime_type: str) -> dict:
        """Parse document and extract text content"""

        # Download file from storage
        file_stream = await self.storage_service.download_file(object_key)

        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(file_stream.read())
            temp_path = temp_file.name

        try:
            if mime_type == "application/pdf":
                return await self.parse_pdf(temp_path)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return await self.parse_docx(temp_path)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return await self.parse_pptx(temp_path)
            elif mime_type == "text/plain":
                return await self.parse_text(temp_path)
            elif mime_type == "application/rtf":
                return await self.parse_rtf(temp_path)
            else:
                raise ValueError(f"Unsupported document type: {mime_type}")

        finally:
            os.unlink(temp_path)

    async def parse_pdf(self, file_path: str) -> dict:
        """Parse PDF document using PyPDF2/pdfplumber"""
        try:
            import PyPDF2

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                pages = []
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    pages.append({
                        "page_number": page_num + 1,
                        "text": text,
                        "word_count": len(text.split()),
                        "char_count": len(text)
                    })

                return {
                    "document_type": "pdf",
                    "total_pages": len(pages),
                    "pages": pages,
                    "metadata": {
                        "encrypted": pdf_reader.is_encrypted,
                        "total_words": sum(p["word_count"] for p in pages),
                        "total_chars": sum(p["char_count"] for p in pages)
                    }
                }
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            # Fallback to simple text extraction
            return {"document_type": "pdf", "pages": [{"page_number": 1, "text": "PDF parsing failed"}]}

    async def parse_docx(self, file_path: str) -> dict:
        """Parse DOCX document using python-docx"""
        try:
            from docx import Document

            doc = Document(file_path)

            paragraphs = []
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    paragraphs.append({
                        "paragraph_number": i + 1,
                        "text": para.text,
                        "style": para.style.name if para.style else "Normal"
                    })

            return {
                "document_type": "docx",
                "paragraphs": paragraphs,
                "metadata": {
                    "total_paragraphs": len(paragraphs),
                    "total_words": sum(len(p["text"].split()) for p in paragraphs),
                    "total_chars": sum(len(p["text"]) for p in paragraphs)
                }
            }
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            return {"document_type": "docx", "paragraphs": [{"text": "DOCX parsing failed"}]}

    async def parse_pptx(self, file_path: str) -> dict:
        """Parse PPTX presentation using python-pptx"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            slides = []
            for slide_num, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_number": slide_num + 1,
                    "shapes": []
                }

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["shapes"].append({
                            "type": str(shape.shape_type),
                            "text": shape.text
                        })

                slides.append(slide_content)

            return {
                "document_type": "pptx",
                "slides": slides,
                "metadata": {
                    "total_slides": len(slides),
                    "total_words": sum(
                        len(shape["text"].split())
                        for slide in slides
                        for shape in slide["shapes"]
                    )
                }
            }
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            return {"document_type": "pptx", "slides": [{"text": "PPTX parsing failed"}]}

    async def parse_text(self, file_path: str) -> dict:
        """Parse plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

                lines = content.split('\n')
                return {
                    "document_type": "text",
                    "content": content,
                    "lines": [{"line_number": i + 1, "text": line} for i, line in enumerate(lines)],
                    "metadata": {
                        "total_lines": len(lines),
                        "total_words": len(content.split()),
                        "total_chars": len(content)
                    }
                }
        except Exception as e:
            logger.error(f"Error parsing text file: {e}")
            return {"document_type": "text", "content": "Text parsing failed"}

    async def parse_rtf(self, file_path: str) -> dict:
        """Parse RTF document"""
        try:
            # Simple RTF parsing - in production, use striprtf or similar
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()

                # Very basic RTF text extraction
                import re
                text = re.sub(r'\\[a-z]+\d*', '', content)  # Remove RTF commands
                text = re.sub(r'[{}]', '', text)  # Remove braces
                text = text.replace('\\', '')

                return {
                    "document_type": "rtf",
                    "content": text.strip(),
                    "metadata": {
                        "total_words": len(text.split()),
                        "total_chars": len(text)
                    }
                }
        except Exception as e:
            logger.error(f"Error parsing RTF: {e}")
            return {"document_type": "rtf", "content": "RTF parsing failed"}

    async def parse_image(self, object_key: str) -> dict:
        """Parse image using OCR (placeholder)"""
        # In production, use PaddleOCR or Tesseract
        return {
            "document_type": "image",
            "ocr_text": "Image OCR not implemented in demo",
            "metadata": {
                "total_words": 0,
                "total_chars": 0
            }
        }

    async def extract_media_metadata(self, object_key: str, mime_type: str) -> dict:
        """Extract metadata from audio/video files"""
        # In production, use ffprobe
        return {
            "media_type": "audio" if mime_type.startswith("audio") else "video",
            "metadata": {
                "duration": 120.0,  # Mock duration
                "format": mime_type
            }
        }

    async def save_parsed_content(self, content: dict, object_key: str):
        """Save parsed content to storage as JSON"""
        import json
        import io

        json_content = json.dumps(content, indent=2).encode('utf-8')
        json_stream = io.BytesIO(json_content)

        # Create a mock UploadFile-like object for storage
        class MockFile:
            def __init__(self, content, filename):
                self.file = io.BytesIO(content)
                self.size = len(content)
                self.content_type = "application/json"
                self.filename = filename

        mock_file = MockFile(json_content, "parsed.json")
        await self.storage_service.upload_file(mock_file, object_key)

async def main():
    """Main entry point for the parser worker"""
    worker = ParserWorker()
    await worker.start()

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    asyncio.run(main())